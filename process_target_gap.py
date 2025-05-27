import os
import json
import traceback
import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# === Google Drive Setup ===
drive_service = None
try:
    creds_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if creds_json:
        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/drive"]
        )
        drive_service = build("drive", "v3", credentials=creds)
except Exception as e:
    print(f"❌ Drive setup failed: {e}")
    traceback.print_exc()

def upload_to_drive(file_path, session_id):
    try:
        query = f"name='{session_id}' and mimeType='application/vnd.google-apps.folder'"
        results = drive_service.files().list(q=query, fields="files(id)").execute()
        folders = results.get("files", [])
        if folders:
            folder_id = folders[0]["id"]
        else:
            folder = drive_service.files().create(body={
                "name": session_id,
                "mimeType": "application/vnd.google-apps.folder"
            }, fields="id").execute()
            folder_id = folder["id"]

        file_meta = {"name": os.path.basename(file_path), "parents": [folder_id]}
        media = MediaFileUpload(file_path, resumable=True)
        uploaded = drive_service.files().create(body=file_meta, media_body=media, fields="id").execute()
        return f"https://drive.google.com/file/d/{uploaded['id']}/view"
    except Exception as e:
        print(f"❌ Upload error: {e}")
        return None

def download_files(files, folder_path):
    downloaded = []
    for f in files:
        if not f.get("file_url"):
            print(f"⚠️ Skipping {f.get('file_name')} – no file_url")
            continue
        try:
            path = os.path.join(folder_path, f["file_name"])
            r = requests.get(f["file_url"], timeout=10)
            with open(path, "wb") as out:
                out.write(r.content)
            f["local_path"] = path
            downloaded.append(f)
        except Exception as e:
            print(f"❌ Download failed for {f['file_name']}: {e}")
    return downloaded
def extract_hw_sw_gaps(files):
    hw, sw = [], []
    for f in files:
        if f["file_type"] in ["gap_hw", "gap_sw"]:
            wb = load_workbook(f["local_path"])
            sheet = wb.active
            for row in sheet.iter_rows(min_row=2, values_only=True):
                platform = str(row[2])
                tier = str(row[3])
                status = str(row[4])
                recommendation = str(row[5]) if row[5] else ""
                entry = {
                    "platform": platform,
                    "tier": tier,
                    "status": status,
                    "recommendation": recommendation
                }
                if f["file_type"] == "gap_hw":
                    hw.append(entry)
                else:
                    sw.append(entry)
    return hw, sw

def create_docx(session_id, folder_path, hw_data, sw_data):
    path = os.path.join(folder_path, f"Target_GAP_Analysis_Report.docx")
    doc = Document()
    doc.add_heading("Target GAP Analysis Report", 0)
    doc.add_paragraph(f"Session: {session_id}\n")

    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph("This document compares the current infrastructure with the target architecture.")

    doc.add_heading("2. GAP Matrix by Domain", level=1)
    doc.add_paragraph("This section presents technical gaps by IT domain (compute, storage, cloud, etc.).")

    def add_table(data, title):
        doc.add_heading(title, level=2)
        table = doc.add_table(rows=1, cols=5)
        hdr = table.rows[0].cells
        hdr[0].text = "Platform"
        hdr[1].text = "Tier"
        hdr[2].text = "Status"
        hdr[3].text = "Recommendation"
        hdr[4].text = "Severity"
        for d in data:
            row = table.add_row().cells
            row[0].text = d["platform"]
            row[1].text = d["tier"]
            row[2].text = d["status"]
            row[3].text = d["recommendation"]
            row[4].text = "5" if "obsolete" in d["status"].lower() else "2"

    add_table(hw_data, "Compute & HW GAPs")
    add_table(sw_data, "Applications & SW GAPs")

    doc.add_heading("3. Functional Impact", level=1)
    doc.add_paragraph("This section outlines how the current gaps impact scalability, uptime, compliance, etc.")

    doc.add_heading("4. Recommendations", level=1)
    for item in hw_data + sw_data:
        if item["recommendation"]:
            doc.add_paragraph(f"- {item['platform']}: {item['recommendation']}")

    doc.add_heading("5. Summary & Observations", level=1)
    doc.add_paragraph("This GAP analysis reveals modernization priorities and transformation areas.")

    doc.save(path)
    return path

def create_pptx(session_id, folder_path, hw_data, sw_data):
    path = os.path.join(folder_path, f"Target_GAP_Analysis_Executive_Report.pptx")
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "Target GAP Executive Report"
    slide.placeholders[1].text = f"Session: {session_id}"

    def add_slide(title, items):
        s = ppt.slides.add_slide(ppt.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.clear()
        for item in items:
            tf.add_paragraph().text = item

    add_slide("HW GAP Summary", [f"{d['platform']} → {d['recommendation']}" for d in hw_data if d["recommendation"]])
    add_slide("SW GAP Summary", [f"{d['platform']} → {d['recommendation']}" for d in sw_data if d["recommendation"]])
    add_slide("Priority Recommendations", [f"{d['platform']}" for d in hw_data + sw_data if "obsolete" in d["status"].lower()])
    ppt.save(path)
    return path
def process_target_gap(session_id, email, files, folder_path):
    try:
        os.makedirs(folder_path, exist_ok=True)

        downloaded = download_files(files, folder_path)
        hw_data, sw_data = extract_hw_sw_gaps(downloaded)

        docx_path = create_docx(session_id, folder_path, hw_data, sw_data)
        pptx_path = create_pptx(session_id, folder_path, hw_data, sw_data)

        docx_url = upload_to_drive(docx_path, session_id)
        pptx_url = upload_to_drive(pptx_path, session_id)

        for f in downloaded:
            f["file_url"] = upload_to_drive(f["local_path"], session_id)

        downloaded.extend([
            {
                "file_name": os.path.basename(docx_path),
                "file_url": docx_url,
                "file_type": "docx_target_gap"
            },
            {
                "file_name": os.path.basename(pptx_path),
                "file_url": pptx_url,
                "file_type": "pptx_target_gap"
            }
        ])

        NEXT_GPT_URL = "https://it-compliance-api.onrender.com/start_it_compliance"
        payload = {
            "session_id": session_id,
            "email": email,
            "gpt_module": "gap_target",
            "files": downloaded,
            "status": "complete"
        }
        requests.post(NEXT_GPT_URL, json=payload)

    except Exception as e:
        print(f"🔥 Target GAP processing failed: {e}")
        traceback.print_exc()
